from typing import List, AsyncGenerator
from io import BytesIO
from pptx import Presentation
from pptx.exc import InvalidXmlError
from tika import parser
from zipfile import BadZipFile
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io
from PIL import Image
import pytesseract
import pybase64
import uuid

from querent.ingestors.ingestor_factory import IngestorFactory
from querent.processors.async_processor import AsyncProcessor
from querent.ingestors.base_ingestor import BaseIngestor
from querent.config.ingestor.ingestor_config import IngestorBackend
from querent.common.types.collected_bytes import CollectedBytes
from querent.common import common_errors
from querent.common.types.ingested_tokens import IngestedTokens
from querent.common.types.ingested_images import IngestedImages
from querent.logging.logger import setup_logger


class PptIngestorFactory(IngestorFactory):
    SUPPORTED_EXTENSIONS = {"ppt", "pptx"}

    async def supports(self, file_extension: str) -> bool:
        return file_extension.lower() in self.SUPPORTED_EXTENSIONS

    async def create(
        self, file_extension: str, processors: List[AsyncProcessor]
    ) -> BaseIngestor:
        if not await self.supports(file_extension):
            return None
        return PptIngestor(processors)


class PptIngestor(BaseIngestor):
    def __init__(self, processors: List[AsyncProcessor]):
        super().__init__(IngestorBackend.PPT)
        self.processors = processors
        self.logger = setup_logger(__name__, "PptIngestor")

    async def ingest(
        self, poll_function: AsyncGenerator[CollectedBytes, None]
    ) -> AsyncGenerator[IngestedTokens, None]:
        current_file = None
        collected_bytes = b""
        try:
            async for chunk_bytes in poll_function:
                if chunk_bytes.is_error() or chunk_bytes.is_eof():
                    continue

                if current_file is None:
                    current_file = chunk_bytes.file
                elif current_file != chunk_bytes.file:
                    async for ingested_data in self.extract_and_process_ppt(
                        CollectedBytes(file=current_file, data=collected_bytes)
                    ):
                        yield ingested_data
                    yield IngestedTokens(
                        file=current_file,
                        data=None,
                        error=None,
                    )
                    collected_bytes = b""
                    current_file = chunk_bytes.file
                collected_bytes += chunk_bytes.data
        except Exception as e:
            yield IngestedTokens(file=current_file, data=None, error=f"Exception: {e}")
        finally:
            async for ingested_data in self.extract_and_process_ppt(
                CollectedBytes(file=current_file, data=collected_bytes)
            ):
                yield ingested_data
            yield IngestedTokens(file=current_file, data=None, error=None)

    async def extract_and_process_ppt(
        self, collected_bytes: CollectedBytes
    ) -> AsyncGenerator[str, None]:
        try:
            if collected_bytes.extension == "pptx":
                ppt_file = BytesIO(collected_bytes.data)
                presentation = Presentation(ppt_file)
                i=1
                for slide in presentation.slides:
                    text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or (shape.shape_type in [MSO_SHAPE_TYPE.PLACEHOLDER, MSO_SHAPE_TYPE.AUTO_SHAPE] and hasattr(shape, "image")):
                            ocr_text = await self.process_image(shape)
                            yield IngestedImages(file=collected_bytes.file, image=pybase64.b64encode(shape.image.blob), image_name=str(uuid.uuid4()), page_num=i, text = text, ocr_text=[ocr_text], coordinates=None)
                    slide_text = "\n".join(text)
                    processed_slide_text = await self.process_data(slide_text)
                    yield IngestedTokens(
                        file=collected_bytes.file, data=processed_slide_text, error=None
                    )
                    i+=1
            elif collected_bytes.extension == "ppt":
                parsed = parser.from_buffer(collected_bytes.data)
                extracted_text = parsed["content"]
                yield IngestedTokens(
                    file=collected_bytes.file, data=extracted_text, error=None
                )
            else:
                raise common_errors.WrongPptFileError(
                    f"Given file is not ppt {collected_bytes.file}"
                )
        except InvalidXmlError as exc:
            raise common_errors.InvalidXmlError(
                f"The following file is not in proper xml format {collected_bytes.file}"
            ) from exc
        except BadZipFile as exc:
            raise common_errors.BadZipFile(
                f"The following file is not a zip file{collected_bytes.file}"
            ) from exc
        except Exception as exc:
            raise common_errors.UnknownError(
                f"Received Unknown error as {exc} from file {collected_bytes.file}"
            ) from exc

    async def process_data(self, text: str) -> str:
        if self.processors is None or len(self.processors) == 0:
            return [text]
        try:
            processed_data = text
            for processor in self.processors:
                processed_data = await processor.process_text(processed_data)
            return processed_data
        except Exception as e:
            self.logger.error(f"Error while processing text: {e}")

    async def process_image(self, shape):
        try:
            # Retrieve image as a BytesIO object
            image_stream = io.BytesIO(shape.image.blob)
            image = Image.open(image_stream)

            # Perform OCR using pytesseract
            ocr_text = pytesseract.image_to_string(image)

            return ocr_text
        except Exception as e:
            self.logger.error(f"Error during image processing: {e}")
            return ""
